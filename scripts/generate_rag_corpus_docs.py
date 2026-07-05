"""Generate realistic enterprise wealth-management / AGP practice-management documents
in PDF, DOCX and PPTX so the RAG corpus exercises the real multi-format parsers, not
just .txt (CLAUDE.md 9.8). Idempotent — overwrites its own outputs.

Uses reportlab (PDF), python-docx (DOCX), python-pptx (PPTX).
"""
from __future__ import annotations

from pathlib import Path

OUT = Path(__file__).resolve().parents[1] / "data" / "documents" / "sample_knowledge"
OUT.mkdir(parents=True, exist_ok=True)


# ---- content (real, substantive practice-management / compliance text) -----------------

PDF_DOCS = {
    "reg_bi_suitability_manual.pdf": (
        "Regulation Best Interest & Suitability Manual",
        [
            ("1. Purpose", "This manual governs how advisors evaluate and document that a recommendation "
             "is in the client's best interest before it is presented. It applies to all managed-account, "
             "advisory and brokerage recommendations firm-wide."),
            ("2. Best-Interest Standard", "At the point of recommendation the advisor must have a reasonable "
             "basis to believe the recommendation is in the retail client's best interest and does not place "
             "the advisor's interest ahead of the client's. Cost, risk, and reasonably available alternatives "
             "must be considered and documented."),
            ("3. Supervisory Review Threshold", "Recommendations with an estimated revenue impact or transfer "
             "value at or above $50,000 require supervisory principal review before presentation. The reviewing "
             "principal must be independent of the recommending advisor's production credit."),
            ("4. Suitability Documentation", "No managed account, advisory program, or discretionary mandate may "
             "be recommended without a current suitability assessment on file, dated within the preceding twelve "
             "months, capturing investment objective, time horizon, liquidity needs and risk tolerance."),
            ("5. Prohibited Representations", "Advisors must not state or imply guaranteed returns, assured "
             "outperformance, or risk-free outcomes in any client communication. Comparative claims against named "
             "competitors require compliance pre-approval."),
            ("6. Record Retention", "Recommendation rationales, suitability assessments and best-interest "
             "determinations are retained for seven years and must be producible within ten business days of a "
             "regulatory or internal audit request."),
        ],
    ),
    "market_outlook_2026_q3.pdf": (
        "Quarterly Market Outlook — 2026 Q3",
        [
            ("Macro Summary", "Growth is moderating as policy rates hold. Base case is a soft landing with "
             "disinflation continuing through year end. Duration risk is balanced; credit spreads remain tight."),
            ("Equities", "Prefer quality and free-cash-flow durability over high-beta cyclicals. Managed equity "
             "sleeves should trim concentrated single-name exposure and rebalance toward the household's agreed "
             "target mix."),
            ("Fixed Income", "Intermediate investment-grade offers the best risk-adjusted carry. Extend duration "
             "opportunistically on rate spikes. Avoid reaching for yield in lower-quality credit."),
            ("Advisor Actions", "Use the outlook to frame proactive review meetings. Households with negative net "
             "cash flow or stale engagement should be prioritized for a planning conversation this quarter."),
        ],
    ),
}

DOCX_DOCS = {
    "managed_account_operating_procedures.docx": (
        "Managed Account Program — Operating Procedures",
        [
            ("Eligibility & Onboarding", "Households are eligible for a managed mandate when investable assets and "
             "suitability support discretionary management. Onboarding requires a signed advisory agreement, a "
             "documented investment policy statement, and a funded account within thirty days."),
            ("Model Assignment", "Each account is mapped to an approved model portfolio aligned to the household's "
             "risk profile. Off-model positions require a documented exception and quarterly review."),
            ("Rebalancing", "Accounts are drift-band rebalanced. Advisors should surface managed-account conversion "
             "opportunities for high-AUM households with low managed penetration and document suitability for the "
             "mandate before presenting."),
            ("Fee Transparency", "At the point of recommendation the client must receive the program fee schedule "
             "and a plain-language statement of how the advisor is compensated. For brokerage-to-advisory transfers "
             "a cost comparison worksheet is mandatory."),
            ("Ongoing Review", "Every advised household receives a documented review at least once every twelve "
             "months capturing current objectives, risk-tolerance confirmation, and any material life changes."),
        ],
    ),
    "agp_milestone_recovery_playbook.docx": (
        "AGP Milestone Recovery Playbook",
        [
            ("When to Trigger", "A recovery plan is triggered when milestone attainment falls below target with "
             "limited days remaining, or when the AGP off-track risk score enters the attention band (40+)."),
            ("Diagnose the Gap", "Decompose the risk score into its drivers: attainment gap, time pressure and CRM "
             "execution risk. Overdue lead and referral follow-ups are the most common and most recoverable driver."),
            ("Recovery Actions", "Work overdue follow-ups oldest-first, refresh the next action on every open "
             "opportunity, advance or close stage-stalled deals, and book a coaching session before the milestone "
             "due date. Recovery plans name specific activities and dates, never generic intentions."),
            ("Escalation", "Persisting off-track status escalates to the district coach. Good-faith self-reported "
             "gaps are treated as mitigating and coached, not penalized."),
        ],
    ),
}

PPTX_DOCS = {
    "agp_onboarding_deck.pptx": (
        "Advisor Growth Program — Onboarding",
        [
            ("What the AGP Is", ["A 24-month structured growth program", "Milestones every three months",
                                 "KPIs across revenue, AUM, leads and referrals", "Coaching tied to measurable impact"]),
            ("How You Are Measured", ["On/off-track risk score, 0-100", "Milestone attainment vs target",
                                       "KPI on-track ratio", "CRM execution quality"]),
            ("What Good Looks Like by Year End", ["Milestone attainment at or above 95%", "KPI on-track ratio above 0.75",
                                                   "NNM positive in at least three of four quarters",
                                                   "Managed-asset share moving toward the agreed mix target"]),
            ("Your Coaching Cadence", ["Quarterly milestone review", "Prioritize the two largest open opportunities",
                                        "Clear overdue follow-ups", "Document next actions in CRM"]),
        ],
    ),
    "referral_coi_workshop.pptx": (
        "Referral & Centre-of-Influence Prospecting Workshop",
        [
            ("Core Principle", ["The warmest path to a new household is an introduction",
                                "From a satisfied client or a centre of influence", "Referral-led prospecting compounds"]),
            ("Play 1 — The Post-Review Referral Ask", ["Ask right after a well-received annual review",
                                                        "Be specific: who in your circle faces a similar decision?",
                                                        "Log resulting names as referral-sourced leads the same day"]),
            ("Play 2 — Centre-of-Influence Mapping", ["List the CPAs and estate attorneys in your book",
                                                       "Draft a value-first introduction", "Track each as an open referral"]),
            ("Measuring Success", ["Referral conversion rate", "Time-to-first-follow-up under 48 hours",
                                   "New households sourced by referral vs cold outreach"]),
        ],
    ),
}


def write_pdfs() -> list[str]:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    styles = getSampleStyleSheet()
    made = []
    for fname, (title, sections) in PDF_DOCS.items():
        path = OUT / fname
        doc = SimpleDocTemplate(str(path), pagesize=LETTER, title=title)
        flow = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
        for heading, body in sections:
            flow.append(Paragraph(heading, styles["Heading2"]))
            flow.append(Paragraph(body, styles["BodyText"]))
            flow.append(Spacer(1, 8))
        doc.build(flow)
        made.append(fname)
    return made


def write_docx() -> list[str]:
    from docx import Document

    made = []
    for fname, (title, sections) in DOCX_DOCS.items():
        d = Document()
        d.add_heading(title, level=0)
        for heading, body in sections:
            d.add_heading(heading, level=1)
            d.add_paragraph(body)
        d.save(str(OUT / fname))
        made.append(fname)
    return made


def write_pptx() -> list[str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    made = []
    for fname, (title, slides) in PPTX_DOCS.items():
        prs = Presentation()
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = title
        cover.placeholders[1].text = "iPerform Practice Management"
        for heading, bullets in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading
            body = slide.placeholders[1].text_frame
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(18)
        prs.save(str(OUT / fname))
        made.append(fname)
    return made


if __name__ == "__main__":
    pdfs = write_pdfs()
    docs = write_docx()
    ppts = write_pptx()
    print("PDF :", pdfs)
    print("DOCX:", docs)
    print("PPTX:", ppts)
    print(f"Wrote {len(pdfs) + len(docs) + len(ppts)} documents to {OUT}")

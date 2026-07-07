"""Dashboard export → PDF / PPTX (real data, server-side render).

Generalizes the presentation-export capability to the whole executive view: pulls the
same real payloads the Executive Dashboard renders (`ScopeDashboardService` +
`AumNetFlowsService`) and lays them into a PDF (reportlab) or a PPTX deck (python-pptx).
No hardcoded numbers — every figure is the live scope+period rollup, so an exported file
matches what the screen shows.
"""

from __future__ import annotations

import io

from app.scope.dashboard import ScopeDashboardService
from app.scope.net_flows import AumNetFlowsService


def _usd(v: float | int | None) -> str:
    if v is None:
        return "—"
    n = float(v)
    sign = "-" if n < 0 else ""
    a = abs(n)
    if a >= 1e9:
        return f"{sign}${a / 1e9:.2f}B"
    if a >= 1e6:
        return f"{sign}${a / 1e6:.2f}M"
    if a >= 1e3:
        return f"{sign}${a / 1e3:.1f}K"
    return f"{sign}${a:,.0f}"


def _gather(scope_type: str, scope_id: str, period: str, compare_to: str) -> dict:
    dash = ScopeDashboardService().dashboard(scope_type, scope_id, period, compare_to)
    flows = AumNetFlowsService().waterfall(scope_type, scope_id, period)
    return {"dash": dash, "flows": flows}


def _kpi_rows(dash: dict) -> list[tuple[str, str]]:
    t = dash.get("totals", {})
    head = dash.get("headline", {})
    return [
        ("Advisors in scope", str(t.get("advisor_count", "—"))),
        (f"Revenue ({dash.get('period')})", _usd(head.get("revenue"))),
        (f"Δ vs {head.get('compare_to', 'Prior')}", f"{head.get('delta_pct', 0):+.1f}%"),
        ("AUM", _usd(t.get("aum_total"))),
        ("NNM (annualized)", _usd(t.get("nnm_annualized"))),
        ("Managed revenue", _usd(t.get("managed_revenue"))),
        ("Avg goal attainment", f"{t.get('avg_goal_attainment', '—')}%"),
        ("Avg AGP risk", str(t.get("avg_agp_risk_score", "—"))),
    ]


def _title(dash: dict) -> str:
    return f"{dash.get('scope_type')} {dash.get('scope_id')} — Executive Summary ({dash.get('period')})"


# --------------------------------------------------------------------------- PDF
def build_pdf(scope_type: str, scope_id: str, period: str, compare_to: str) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    )
    from reportlab.lib.styles import getSampleStyleSheet

    data = _gather(scope_type, scope_id, period, compare_to)
    dash, flows = data["dash"], data["flows"]
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.6 * inch)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("iPerform Insights and Coaching", styles["Heading3"]))
    story.append(Paragraph(_title(dash), styles["Title"]))
    story.append(Paragraph(f"Compare to: {dash.get('compare_to')}", styles["Normal"]))
    story.append(Spacer(1, 0.2 * inch))

    def _table(rows, header=None, col_widths=None):
        body = ([header] if header else []) + rows
        tbl = Table(body, colWidths=col_widths, hAlign="LEFT")
        style = [
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#E2E8F0")),
            ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
        ]
        if header:
            style += [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F172A")),
                      ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                      ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold")]
        tbl.setStyle(TableStyle(style))
        return tbl

    story.append(Paragraph("Key Metrics", styles["Heading2"]))
    story.append(_table(_kpi_rows(dash), col_widths=[2.6 * inch, 2.2 * inch]))
    story.append(Spacer(1, 0.2 * inch))

    if flows.get("available"):
        story.append(Paragraph("AUM Net-Flows Bridge", styles["Heading2"]))
        story.append(_table(
            [[s["label"], _usd(s["value"])] for s in flows["steps"]],
            header=["Component", "Amount"], col_widths=[2.6 * inch, 2.2 * inch]))
        story.append(Spacer(1, 0.2 * inch))

    for key, heading in (("top_advisors", "Top Advisors"), ("bottom_advisors", "Bottom Advisors")):
        rows = dash.get(key, [])[:5]
        if rows:
            story.append(Paragraph(heading, styles["Heading2"]))
            story.append(_table(
                [[r.get("advisor_name", r.get("advisor_id")), _usd(r.get("revenue_ltm")),
                  _usd(r.get("aum_total")), (r.get("reason") or "")[:48]] for r in rows],
                header=["Advisor", "Revenue", "AUM", "Why"],
                col_widths=[1.5 * inch, 1.1 * inch, 1.1 * inch, 2.4 * inch]))
            story.append(Spacer(1, 0.15 * inch))

    story.append(Spacer(1, 0.1 * inch))
    story.append(Paragraph(
        "Every figure is a live rollup over real per-advisor snapshots and transactions.",
        styles["Italic"]))
    doc.build(story)
    return buf.getvalue()


# -------------------------------------------------------------------------- PPTX
def build_pptx(scope_type: str, scope_id: str, period: str, compare_to: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    data = _gather(scope_type, scope_id, period, compare_to)
    dash, flows = data["dash"], data["flows"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    navy = RGBColor(0x0F, 0x17, 0x2A)

    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = _title(dash)
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = f"iPerform Insights and Coaching · Compare to {dash.get('compare_to')}"
    p.font.size = Pt(18)
    p.font.color.rgb = navy

    def _table_slide(heading, header, rows):
        sl = prs.slides.add_slide(prs.slide_layouts[5])
        sl.shapes.title.text = heading
        n = len(rows) + 1
        cols = len(header)
        tbl = sl.shapes.add_table(n, cols, Inches(0.5), Inches(1.5),
                                  Inches(12.3), Inches(0.4 * n)).table
        for c, h in enumerate(header):
            cell = tbl.cell(0, c)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(11)

    _table_slide("Key Metrics", ["Metric", "Value"], _kpi_rows(dash))
    if flows.get("available"):
        _table_slide("AUM Net-Flows Bridge", ["Component", "Amount"],
                     [[s2["label"], _usd(s2["value"])] for s2 in flows["steps"]])
    for key, heading in (("top_advisors", "Top Advisors"), ("bottom_advisors", "Bottom Advisors")):
        rows = dash.get(key, [])[:6]
        if rows:
            _table_slide(heading, ["Advisor", "Revenue", "AUM", "Why"],
                         [[r.get("advisor_name", r.get("advisor_id")), _usd(r.get("revenue_ltm")),
                           _usd(r.get("aum_total")), (r.get("reason") or "")[:60]] for r in rows])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

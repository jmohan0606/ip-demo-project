from __future__ import annotations

from pathlib import Path


class DocumentParser:
    """Extracts plain text from knowledge documents for chunking + embedding.

    Real extraction for text, PDF (pypdf), DOCX (python-docx) and PPTX (python-pptx).
    No OCR — scanned/image-only PDFs yield little/no text and should be flagged upstream.
    """

    SUPPORTED_TEXT_SUFFIXES = {".txt", ".md", ".csv", ".json", ".html"}

    def parse(self, path: str | Path) -> str:
        file_path = Path(path)
        if not file_path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")

        suffix = file_path.suffix.lower()
        if suffix in self.SUPPORTED_TEXT_SUFFIXES:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            return self._parse_pdf(file_path)
        if suffix == ".docx":
            return self._parse_docx(file_path)
        if suffix == ".pptx":
            return self._parse_pptx(file_path)

        raise ValueError(f"Unsupported document suffix: {suffix}")

    @staticmethod
    def _parse_pdf(file_path: Path) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
        return "\n\n".join(p for p in pages if p)

    @staticmethod
    def _parse_docx(file_path: Path) -> str:
        from docx import Document

        document = Document(str(file_path))
        parts = [para.text for para in document.paragraphs if para.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    @staticmethod
    def _parse_pptx(file_path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(str(file_path))
        parts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ]
            if slide_texts:
                parts.append(f"[Slide {index}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
